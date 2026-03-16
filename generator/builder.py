from pptx import Presentation
from pptx.util import Inches


# --------------------------------------------------
# SAFE TITLE SETTER
# --------------------------------------------------
def set_slide_title(slide, text):

    if slide.shapes.title is not None:
        slide.shapes.title.text = text
        return

    textbox = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(0.3),
        Inches(9),
        Inches(1)
    )

    textbox.text_frame.text = text


# --------------------------------------------------
# SAFE BODY WRITER
# --------------------------------------------------
def write_body(slide, content):

    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = content
        return

    textbox = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(1.5),
        Inches(9),
        Inches(4)
    )

    textbox.text_frame.text = content

def group_by_column(slide):

    elements = slide["elements"]

    # sort by vertical position first
    elements = sorted(elements, key=lambda x: (x["left"]))

    # detect header row (top-most unique X)
    headers = {}

    threshold_x = 500000
    threshold_y = 800000

    for el in elements:

        placed = False

        for h in headers:

            if abs(h - el["left"]) < threshold_x:
                headers[h].append(el)
                placed = True
                break

        if not placed:
            headers[el["left"]] = [el]

    groups = []

    for col in headers.values():

        col = sorted(col, key=lambda x: x.get("top", 0))

        header = col[0]["text"]

        body = [e["text"] for e in col[1:]]

        groups.append([header] + body)

    return groups

# --------------------------------------------------
# FIND BEST LAYOUT BASED ON PLACEHOLDER COUNT
# --------------------------------------------------
def find_best_layout(prs, num_groups):

    best_layout = prs.slide_layouts[0]
    best_diff = 999

    for layout in prs.slide_layouts:

        capacity = len(layout.placeholders)

        diff = abs(capacity - (num_groups + 1))  # +1 for title

        if diff < best_diff:
            best_diff = diff
            best_layout = layout

    return best_layout


# --------------------------------------------------
# MAIN BUILDER
# --------------------------------------------------
def build(template, slides, ai_plan):

    prs = Presentation(template)

    for slide_json, plan in zip(slides, ai_plan):

        groups = group_by_column(slide_json)

        # ------------------------------------------
        # GRID → SPLIT INTO MULTIPLE SLIDES
        # ------------------------------------------
        if plan["split_required"]:

            for g in groups:

                layout = find_best_layout(prs, 1)

                slide = prs.slides.add_slide(layout)

                set_slide_title(slide, g[0][:50])

                write_body(slide, "\n".join(g))

        # ------------------------------------------
        # NORMAL → SINGLE STRUCTURED SLIDE
        # ------------------------------------------
        else:

            layout = find_best_layout(prs, len(groups))

            slide = prs.slides.add_slide(layout)

            set_slide_title(slide, slide_json["title"])

            for i, g in enumerate(groups):

                content = "\n".join(g)

                if i + 1 < len(slide.placeholders):
                    slide.placeholders[i + 1].text = content
                else:
                    textbox = slide.shapes.add_textbox(
                        Inches(0.5),
                        Inches(1.5 + i),
                        Inches(9),
                        Inches(1.2)
                    )
                    textbox.text_frame.text = content

    return prs



#### not use full



# from pptx import Presentation
# from pptx.enum.shapes import PP_PLACEHOLDER


# def group_by_column(slide):

#     cols = {}

#     threshold = 500000

#     for el in slide["elements"]:

#         placed = False

#         for k in cols:
#             if abs(k - el["left"]) < threshold:
#                 cols[k].append(el)
#                 placed = True
#                 break

#         if not placed:
#             cols[el["left"]] = [el]

#     groups = []

#     for k in sorted(cols.keys()):

#         column = sorted(cols[k], key=lambda x: x["top"])

#         header = column[0]["text"]

#         body = [c["text"] for c in column[1:]]

#         groups.append({
#             "header": header,
#             "body": body
#         })

#     return groups


# def find_title_content_layout(prs):

#     for layout in prs.slide_layouts:

#         has_title = False
#         has_body = False

#         for ph in layout.placeholders:

#             if ph.placeholder_format.type == PP_PLACEHOLDER.TITLE:
#                 has_title = True

#             if ph.placeholder_format.type == PP_PLACEHOLDER.BODY:
#                 has_body = True

#         if has_title and has_body:
#             return layout

#     return prs.slide_layouts[1]


# def build(template, slides, ai_plan):

#     prs = Presentation(template)

#     for slide_json, plan in zip(slides, ai_plan):

#         groups = group_by_column(slide_json)

#         layout = find_title_content_layout(prs)

#         for g in groups:

#             slide = prs.slides.add_slide(layout)

#             for ph in slide.placeholders:

#                 if ph.placeholder_format.type == PP_PLACEHOLDER.TITLE:
#                     ph.text = g["header"]

#                 if ph.placeholder_format.type == PP_PLACEHOLDER.BODY:
#                     ph.text = "\n".join(g["body"])

#     return prs