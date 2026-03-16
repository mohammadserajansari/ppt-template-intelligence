from pptx import Presentation

def get_theme_features(template):

    prs = Presentation(template)

    fonts = set()

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            fonts.add(run.font.name)

    return {"fonts": list(fonts)}