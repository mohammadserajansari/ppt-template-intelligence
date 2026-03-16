from pptx import Presentation
from core.logger import logger

def extract_ppt(path):

    logger.info("Extraction started")

    prs = Presentation(path)

    slides = []

    for i, slide in enumerate(prs.slides):

        data = {
            "slide_index": i,
            "title": "",
            "elements": [],
            "tables": []
        }

        for shape in slide.shapes:

            if shape == slide.shapes.title:
                data["title"] = shape.text.strip()
                continue

            if shape.has_text_frame:

                for p in shape.text_frame.paragraphs:

                    txt = p.text.strip()

                    if txt:
                        data["elements"].append({
                            "text": txt,
                            "level": p.level,
                            "left": int(shape.left),
                            "top": int(shape.top)
                        })

            if shape.has_table:

                table = []
                for r in shape.table.rows:
                    table.append([c.text for c in r.cells])

                data["tables"].append(table)

        slides.append(data)

    logger.info("Extraction finished")

    return slides