from pptx import Presentation

def analyze_template(path):

    prs = Presentation(path)

    meta = []

    for l in prs.slide_layouts:
        meta.append({
            "layout": l.name,
            "placeholders": len(l.placeholders)
        })

    return meta